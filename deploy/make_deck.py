"""발표용 제품 소개서(.pptx)를 만든다.

    py -3.13 deploy/make_deck.py        # → Cognitive-Economy-소개서.pptx

내용을 고치고 싶으면 이 파일의 문자열을 고치고 다시 돌리면 된다.
다이어그램도 그림이 아니라 **파워포인트 도형**으로 넣으므로, 만든 뒤
파워포인트에서 그대로 끌어 옮기고 색을 바꿀 수 있다.
"""
from __future__ import annotations

import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

OUT = pathlib.Path(__file__).resolve().parent.parent / "Cognitive-Economy-소개서.pptx"

# ── 색 ────────────────────────────────────────────────────────────
INK = RGBColor(0x0B, 0x1B, 0x34)      # 본문 먹색
DIM = RGBColor(0x5A, 0x6B, 0x82)      # 보조 설명
LINE = RGBColor(0xD8, 0xDF, 0xE8)
BLUE = RGBColor(0x1B, 0x5F, 0xD9)     # 강조
GREEN = RGBColor(0x0E, 0x8A, 0x4F)
RED = RGBColor(0xC7, 0x33, 0x2E)
VIOLET = RGBColor(0x6B, 0x3F, 0xC4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TINT = RGBColor(0xF3, 0xF6, 0xFB)     # 카드 배경
FONT = "맑은 고딕"

W, H = Emu(12192000), Emu(6858000)    # 16:9
M = Emu(720000)                       # 좌우 여백


def cm(v: float) -> Emu:
    return Emu(int(v * 360000))


def deck() -> Presentation:
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def text(slide, x, y, w, h, s, *, size=14, bold=False, color=INK,
         align=PP_ALIGN.LEFT, space=1.35, anchor=MSO_ANCHOR.TOP):
    """문단 여러 개를 한 번에. 줄바꿈이 곧 문단이다."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(str(s).split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = space
        run = para.add_run()
        run.text = line
        f = run.font
        f.name, f.size, f.bold, f.color.rgb = FONT, Pt(size), bold, color
    return box


def box(slide, x, y, w, h, *, fill=TINT, edge=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if edge is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = edge
        shape.line.width = Pt(1.25)
    shape.shadow.inherit = False
    if radius:
        try:
            shape.adjustments[0] = 0.08
        except (IndexError, AttributeError):
            pass
    return shape


def label(shape, s, *, size=12, bold=True, color=INK, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Emu(90000)
    for i, line in enumerate(s.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = 1.2
        run = para.add_run()
        run.text = line
        f = run.font
        f.name, f.size, f.color.rgb = FONT, Pt(size), color
        # 첫 줄만 굵게 — 아래는 설명이라 톤을 낮춘다
        f.bold = bold if i == 0 else False
        if i:
            f.size = Pt(size - 2)
            f.color.rgb = DIM if color == INK else color


def arrow(slide, x1, y1, x2, y2, *, color=BLUE, width=1.5, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dashed:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    c.line.end_arrowhead = True          # python-pptx 미지원 시 무시됨
    return c


def head(slide, title, sub=""):
    """모든 내지의 머리. 위치를 한 곳에서 정해 슬라이드마다 흔들리지 않게."""
    text(slide, M, cm(1.5), W - 2 * M, cm(1.6), title, size=38, bold=True)
    if sub:
        text(slide, M, cm(3.55), W - 2 * M, cm(1.0), sub, size=16, color=DIM)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, cm(4.35),
                                  Emu(int(cm(2.6))), Emu(50000))
    rule.fill.solid()
    rule.fill.fore_color.rgb = BLUE
    rule.line.fill.background()
    rule.shadow.inherit = False


def footer(slide, n):
    text(slide, W - M - cm(4), H - cm(1.3), cm(4), cm(0.7),
         f"Cognitive Economy   {n}", size=11, color=DIM, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
def build() -> None:
    prs = deck()

    # ── 1. 표지 ────────────────────────────────────────────────
    s = blank(prs)
    bg = box(s, 0, 0, W, H, fill=RGBColor(0x0B, 0x1B, 0x34), radius=False)
    bg.line.fill.background()
    text(s, M, cm(4.6), W - 2 * M, cm(1.2),
         "Google Cloud × Solana AI Agentic Hackathon · Track A",
         size=16, color=RGBColor(0x8F, 0xA6, 0xC4))
    text(s, M, cm(6.2), W - 2 * M, cm(3.4), "Cognitive Economy",
         size=76, bold=True, color=WHITE)
    text(s, M, cm(10.2), W - 2 * M, cm(1.6),
         "스스로 벌어서 스스로 쓰는 투자 에이전트",
         size=30, color=RGBColor(0x6E, 0xB4, 0xFF))
    text(s, M, cm(12.8), W - 2 * M, cm(2.4),
         "사람은 처음 한 번만 서명합니다. 그 뒤의 결제·판단·매매는\n"
         "에이전트가 단독으로 집행하고, 쓴 비용을 스스로 벌어 충당합니다.",
         size=18, color=RGBColor(0xB6, 0xC6, 0xDC))

    # ── 2. 한 줄 정의 ──────────────────────────────────────────
    s = blank(prs)
    head(s, "무엇을 만들었나", "실제로 도는 제품입니다 — 시연 링크가 있습니다")
    cards = [
        ("에이전트가 직접 결제합니다", "뉴스·추론을 호출 한 건마다\n온체인으로 결제합니다.\n구독이 아니라 사용량 과금입니다.", BLUE),
        ("사람의 규칙이 최종 거부권", "모델이 아무리 확신해도\n사용자가 정한 하한 미만이면\n체결되지 않습니다.", GREEN),
        ("모든 판단에 영수증", "어떤 뉴스·어떤 모델이 답했는지\n남습니다. 폴백이 일어난 판단은\n팔 수 없습니다.", VIOLET),
    ]
    x = M
    cw = Emu(int((W - 2 * M - cm(0.8) * 2) / 3))
    for title, body, color in cards:
        c = box(s, x, cm(5.4), cw, cm(8.2), edge=LINE)
        bar = box(s, x, cm(5.4), cw, cm(0.22), fill=color, radius=False)
        bar.line.fill.background()
        text(s, x + cm(0.6), cm(6.4), cw - cm(1.2), cm(1.8), title,
             size=18, bold=True, color=color)
        text(s, x + cm(0.6), cm(8.6), cw - cm(1.2), cm(4.4), body,
             size=15, color=INK)
        x += cw + cm(0.8)
    text(s, M, cm(14.6), W - 2 * M, cm(1.6),
         "지금 상태 :  Solana devnet 실거래 · Google Gemini 실추론 · "
         "한국투자증권 실시세 · 4개국 76종목",
         size=17, bold=True, color=BLUE)
    footer(s, 2)

    # ── 3. 문제 ────────────────────────────────────────────────
    s = blank(prs)
    head(s, "무엇이 막고 있나", "AI 에이전트에게 투자 판단을 맡기려 할 때 반드시 만나는 세 벽")
    probs = [
        ("01", "에이전트에게는 지갑이 없습니다",
         "구독과 카드는 사람 명의이고 사람의 서명을 요구합니다. 에이전트가 데이터를 하나 사려면\n"
         "그때마다 사람을 불러야 합니다. 자율이라 부르기 어렵습니다."),
        ("02", "무엇을 근거로 판단했는지 남지 않습니다",
         "어떤 기사를 읽고 어떤 모델이 답했는지 증빙이 없습니다. 금융권이 AI 판단을 도입하지\n"
         "못하는 실질적 이유가 성능이 아니라 여기에 있습니다."),
        ("03", "권한을 주면 한도가 사라집니다",
         "권한을 넓히면 통제가 사라지고, 좁히면 자동화가 끊깁니다. 그 사이를 다룰 장치가\n"
         "없어서 결국 사람이 매번 확인하는 구조로 되돌아갑니다."),
    ]
    y = cm(5.4)
    for num, title, body in probs:
        n = box(s, M, y, cm(2.1), cm(2.1), fill=RGBColor(0xFD, 0xEC, 0xEB))
        label(n, num, size=26, color=RED)
        text(s, M + cm(2.9), y + cm(0.05), W - 2 * M - cm(3.2), cm(1.1),
             title, size=23, bold=True)
        text(s, M + cm(2.9), y + cm(1.45), W - 2 * M - cm(3.2), cm(2.0),
             body, size=16, color=DIM)
        y += cm(3.5)
    footer(s, 3)

    # ── 4. 해결 ────────────────────────────────────────────────
    s = blank(prs)
    head(s, "네 가지 장치로 풉니다", "각각이 위 문제 하나씩을 정면으로 맡습니다")
    pillars = [
        ("x402", "호출당 온체인 결제", "402 챌린지 → 결제 → 증빙 첨부 재요청.\n"
                                    "증빙은 리소스에 묶이고 한 번만 쓰입니다.", BLUE),
        ("룰북", "사람의 최종 거부권", "확신도 하한·1회 최대금액·일일 횟수.\n"
                                    "모델 위에 사람이 있습니다.", GREEN),
        ("영수증", "판단의 증빙", "어떤 뉴스·어떤 모델이 답했는지 기록.\n"
                                "폴백이면 그 판단은 팔 수 없습니다.", VIOLET),
        ("위임", "한도 있는 자율", "SPL approve 한 번. 이후 서명 없이\n"
                                "집행되고, 한도 초과는 체인이 거부합니다.", RED),
    ]
    cw = Emu(int((W - 2 * M - cm(0.6) * 3) / 4))
    x = M
    for tag, title, body, color in pillars:
        box(s, x, cm(5.4), cw, cm(8.4), edge=LINE)
        chip = box(s, x + cm(0.5), cm(6.1), cm(2.6), cm(1.1), fill=color)
        label(chip, tag, size=15, color=WHITE)
        text(s, x + cm(0.5), cm(7.8), cw - cm(1.0), cm(1.3), title,
             size=17, bold=True)
        text(s, x + cm(0.5), cm(9.7), cw - cm(1.0), cm(3.6), body,
             size=13.5, color=DIM)
        x += cw + cm(0.6)
    text(s, M, cm(14.8), W - 2 * M, cm(1.6),
         "네 장치가 함께 있을 때만 성립합니다. 결제만 있으면 통제가 없고, "
         "규칙만 있으면 자동화가 아닙니다.",
         size=17, color=INK)
    footer(s, 4)

    # ── 5. 아키텍처 ────────────────────────────────────────────
    s = blank(prs)
    head(s, "구조", "봇 하나가 지갑 네 개를 갖고, 오갈 수 있는 경로가 고정돼 있습니다")

    # 왼쪽: 사용자
    u = box(s, M, cm(7.4), cm(5.6), cm(2.6), fill=RGBColor(0xE8, 0xF0, 0xFE),
            edge=BLUE)
    label(u, "사용자 지갑\n서명은 여기서 딱 한 번", size=17, color=BLUE)
    text(s, M, cm(10.3), cm(5.6), cm(1.0), "SPL approve\n(한도 위임)",
         size=13, color=DIM, align=PP_ALIGN.CENTER)

    # 가운데: 봇 내부 지갑 4개
    bx = M + cm(7.2)
    bw = cm(13.4)
    box(s, bx, cm(5.5), bw, cm(10.4), fill=WHITE, edge=LINE)
    text(s, bx, cm(5.8), bw, cm(0.8), "봇 — 역할별 지갑 4개 · 경로 화이트리스트",
         size=14, bold=True, color=DIM, align=PP_ALIGN.CENTER)
    wallets = [
        ("user-treasury", "사용자 원금", GREEN),
        ("invest-wallet", "매매 집행", BLUE),
        ("research-agent", "데이터 구매", VIOLET),
        ("revenue-wallet", "수익 적립", RED),
    ]
    ww = Emu(int((bw - cm(1.4) - cm(0.5)) / 2))
    for i, (name, desc, color) in enumerate(wallets):
        wx = bx + cm(0.7) + (ww + cm(0.5)) * (i % 2)
        wy = cm(7.1) + cm(2.4) * (i // 2)
        w = box(s, wx, wy, ww, cm(2.0), fill=TINT, edge=color)
        label(w, f"{name}\n{desc}", size=15, color=color)

    warn = box(s, bx + cm(0.7), cm(12.2), bw - cm(1.4), cm(2.9),
               fill=RGBColor(0xFD, 0xEC, 0xEB))
    label(warn, "user-treasury → research-agent 경로는\n존재하지 않습니다\n"
                "사용자 원금이 API 비용으로 새는 일이 구조적으로 불가능합니다",
          size=15, color=RED)

    # 오른쪽: 외부
    ex = bx + bw + cm(1.0)
    exw = W - M - ex
    ext = [("Google Gemini", "매수·매도 판단", VIOLET),
           ("한국투자증권 OpenAPI", "국내·해외 실시간 시세", BLUE),
           ("Alpha Vantage", "당일 뉴스", GREEN),
           ("Solana devnet", "결제·정산·미러 토큰", RED)]
    for i, (name, desc, color) in enumerate(ext):
        e = box(s, ex, cm(6.0) + cm(2.5) * i, exw, cm(2.0), fill=WHITE,
                edge=color)
        label(e, f"{name}\n{desc}", size=14, color=color)

    arrow(s, M + cm(5.6), cm(8.7), bx, cm(8.7), color=BLUE, width=2.5)
    arrow(s, bx + bw, cm(10.5), ex, cm(10.5), color=DIM, width=2, dashed=True)
    text(s, bx + bw + cm(0.05), cm(9.9), cm(1.4), cm(0.6), "x402",
         size=12, bold=True, color=DIM)
    footer(s, 5)

    # ── 6. 파이프라인 ──────────────────────────────────────────
    s = blank(prs)
    head(s, "한 번의 판단이 지나가는 길", "각 단계가 스스로 결제합니다 — 사람의 서명은 없습니다")
    steps = [
        ("뉴스 구매", "$0.0020", BLUE),
        ("1차 스크리닝", "$0.00015", BLUE),
        ("심층 추론", "$0.0056", VIOLET),
        ("룰북 검사", "무료 · 거부권", GREEN),
        ("체결", "온체인", BLUE),
        ("정산 85/10/5", "한 트랜잭션", RED),
    ]
    sw = Emu(int((W - 2 * M - cm(0.45) * 5) / 6))
    x = M
    for i, (name, cost, color) in enumerate(steps):
        b = box(s, x, cm(6.2), sw, cm(4.4), fill=TINT, edge=color)
        label(b, f"{name}\n{cost}", size=16, color=color)
        if i < 5:
            arrow(s, x + sw, cm(8.4), x + sw + cm(0.45), cm(8.4),
                  color=DIM, width=1.75)
        x += sw + cm(0.45)
    text(s, M, cm(12.0), W - 2 * M, cm(3.8),
         "룰북 검사에서 막히면 앞의 결제는 이미 일어난 뒤입니다. 그래도 되돌리지 않습니다 —\n"
         "판단에 든 비용은 결과와 무관하게 발생한 비용이고, 그 사실을 감추지 않는 것이\n"
         "이 구조의 핵심입니다. 적중률이 40% 아래로 떨어지면 정책이 추가 투입을 거절합니다.",
         size=17, color=INK)
    footer(s, 6)

    # ── 7. 타깃 ────────────────────────────────────────────────
    s = blank(prs)
    head(s, "누구에게 파는가", "세 층이 서로를 필요로 합니다 — 한 층만으로는 시장이 열리지 않습니다")
    targets = [
        ("1차", "증권사 · 은행 리테일 앱",
         "신한 · 토스 · 카카오페이증권 · 미래에셋",
         "자체 AI 자문 기능을 갖고 싶지만 '무엇을 근거로 판단했는가'를\n"
         "설명하지 못해 출시를 못 합니다. 영수증 구조가 그 빈칸을 채웁니다.", BLUE),
        ("2차", "데이터 · 리서치 공급자",
         "증권사 리서치 · 뉴스사 · 시세 API 사업자",
         "구독은 사람 수만큼만 팔립니다. 에이전트에게는 호출당 과금이\n"
         "맞고, 그 창구가 없어서 시장이 안 열려 있습니다.", GREEN),
        ("3차", "AI 에이전트 개발사",
         "핀테크 스타트업 · 자산관리 SaaS",
         "에이전트에 '지출 권한'을 붙이는 일을 각자 다시 만들고 있습니다.\n"
         "결제·한도·감사를 인프라로 가져다 씁니다.", VIOLET),
    ]
    y = cm(5.4)
    for tag, who, examples, why, color in targets:
        chip = box(s, M, y + cm(0.5), cm(2.0), cm(1.1), fill=color)
        label(chip, tag, size=15, color=WHITE)
        text(s, M + cm(2.6), y, cm(10.4), cm(1.1), who, size=22, bold=True)
        text(s, M + cm(2.6), y + cm(1.3), cm(10.4), cm(0.9), examples,
             size=14, color=color)
        text(s, M + cm(14.0), y + cm(0.15), W - M - (M + cm(14.0)), cm(2.4),
             why, size=15, color=DIM)
        y += cm(3.4)
    footer(s, 7)

    # ── 8. 수익 모델 ───────────────────────────────────────────
    s = blank(prs)
    head(s, "어떻게 버는가", "에이전트가 쓰는 돈이 늘수록 매출이 늘어나는 구조입니다")
    revs = [
        ("결제 라우팅 수수료", "거래액의 20~50bp",
         "에이전트가 데이터를 살 때마다\n결제 경로에서 발생합니다.\n호출 수에 정비례합니다."),
        ("데이터 마켓 중개", "거래액의 5~15%",
         "공급자와 에이전트를 잇는\n창구 수수료입니다.\n양쪽이 늘수록 커집니다."),
        ("성과 분배", "정산의 10%",
         "매매 이익이 났을 때만\n가져갑니다(85/10/5).\n고객과 이해가 일치합니다."),
        ("화이트라벨 라이선스", "연간 고정 + MAU",
         "금융사 앱에 내장하는 형태.\n초기 매출의 중심이 됩니다."),
    ]
    cw = Emu(int((W - 2 * M - cm(0.6) * 3) / 4))
    x = M
    for name, price, body in revs:
        box(s, x, cm(5.4), cw, cm(7.4), edge=LINE)
        text(s, x + cm(0.5), cm(6.2), cw - cm(1.0), cm(1.4), name,
             size=16, bold=True)
        text(s, x + cm(0.5), cm(8.0), cw - cm(1.0), cm(1.1), price,
             size=16, bold=True, color=BLUE)
        text(s, x + cm(0.5), cm(9.7), cw - cm(1.0), cm(2.8), body,
             size=13.5, color=DIM)
        x += cw + cm(0.6)
    note = box(s, M, cm(13.8), W - 2 * M, cm(2.6), fill=TINT)
    label(note, "구독 모델과 다른 점 :  사람 수가 아니라 '판단 횟수'에 과금합니다.\n"
                "에이전트는 사람보다 훨씬 자주 판단하므로, 같은 고객사에서 매출이 다르게 자랍니다.",
          size=17, color=INK)
    footer(s, 8)

    # ── 9. 도입 시나리오 ───────────────────────────────────────
    s = blank(prs)
    head(s, "어떻게 들어가는가", "증권사 리테일 앱을 기준으로 한 3단계")
    phases = [
        ("1단계", "0 – 3개월", "파일럿 · 모의 자산",
         "앱 안에 탭 하나를 붙입니다. 미러 자산과 모의 계좌로\n"
         "동작하고, 실제 주문은 내지 않습니다. 고객은 AI가 어떤\n"
         "근거로 판단했는지를 처음으로 보게 됩니다.", GREEN),
        ("2단계", "3 – 9개월", "자사 리서치를 판매",
         "증권사의 리포트·시황을 에이전트가 호출당 구매하도록\n"
         "엽니다. 리서치가 비용 부서에서 매출 부서로 바뀝니다.\n"
         "이 단계에서 데이터 마켓 수수료가 발생합니다.", BLUE),
        ("3단계", "9개월 –", "실계좌 주문 연동",
         "규제 검토를 마친 범위에서 실제 주문을 라우팅합니다.\n"
         "위임 한도와 룰북이 그대로 규제 대응 장치가 됩니다.\n"
         "성과 분배 매출이 여기서 시작됩니다.", VIOLET),
    ]
    cw = Emu(int((W - 2 * M - cm(0.8) * 2) / 3))
    x = M
    for tag, when, title, body, color in phases:
        box(s, x, cm(5.4), cw, cm(8.8), edge=LINE)
        bar = box(s, x, cm(5.4), cw, cm(0.22), fill=color, radius=False)
        bar.line.fill.background()
        text(s, x + cm(0.8), cm(6.2), cw - cm(1.6), cm(0.9),
             f"{tag}   {when}", size=15, bold=True, color=color)
        text(s, x + cm(0.8), cm(7.5), cw - cm(1.6), cm(1.3), title,
             size=21, bold=True)
        text(s, x + cm(0.8), cm(9.4), cw - cm(1.6), cm(4.2), body,
             size=15, color=DIM)
        x += cw + cm(0.8)
    text(s, M, cm(15.0), W - 2 * M, cm(1.4),
         "미러 주식으로 시작하는 이유 :  실제 체결은 자본시장법 정리가 필요한 영역이라, "
         "규제가 열리는 속도에 맞춰 단계를 나눴습니다.",
         size=16, color=DIM)
    footer(s, 9)

    # ── 10. 지금 되는 것 ───────────────────────────────────────
    s = blank(prs)
    head(s, "지금 실제로 되는 것", "슬라이드가 아니라 도는 서버에서 측정한 값입니다")
    facts = [
        ("76", "미러 주식 토큰", "코스피 19 · 코스닥 20\n나스닥 19 · 도쿄 18"),
        ("1 : 1", "호출당 온체인 결제", "API 호출 한 건 =\n트랜잭션 한 건"),
        ("25 / 25", "검증 항목 통과", "재사용 거부 · 원자적 정산\n재시작 생존 포함"),
        ("9 / 9", "보안 감사", "자금 경로 · 권한 · 격리"),
    ]
    cw = Emu(int((W - 2 * M - cm(0.6) * 3) / 4))
    x = M
    for big, title, body in facts:
        box(s, x, cm(5.4), cw, cm(5.6), fill=TINT)
        text(s, x, cm(6.2), cw, cm(2.2), big, size=42, bold=True, color=BLUE,
             align=PP_ALIGN.CENTER)
        text(s, x, cm(8.5), cw, cm(0.9), title, size=17, bold=True,
             align=PP_ALIGN.CENTER)
        text(s, x, cm(9.5), cw, cm(1.6), body, size=14, color=DIM,
             align=PP_ALIGN.CENTER)
        x += cw + cm(0.6)
    text(s, M, cm(11.8), W - 2 * M, cm(4.2),
         "직접 확인하실 수 있습니다\n"
         "· 거래 내역의 서명을 누르면 Solana Explorer 가 열리고, 화면 수량과 "
         "온체인 수량이 소수점까지 같습니다\n"
         "· API 탭에서 호출 한 건마다 결제 한 건, 서명 하나가 남는 것을 보실 수 있습니다\n"
         "· 봇에게 주식과 무관한 질문을 하면 정중히 거절합니다",
         size=16, color=INK)
    footer(s, 10)

    # ── 11. 리스크 통제 ────────────────────────────────────────
    s = blank(prs)
    head(s, "무엇을 막아 두었나", "자율성을 주장하려면 통제부터 증명해야 한다고 보았습니다")
    rows = [
        ("자금 경로", "user-treasury → research-agent 경로가 없습니다. "
                    "사용자 원금이 API 비용으로 새지 않습니다."),
        ("결제 증빙", "리소스에 묶이고 한 번만 쓰입니다. 재사용·교차 사용은 거부됩니다."),
        ("정산 원자성", "85/10/5 세 갈래가 하나의 온체인 트랜잭션을 공유합니다. "
                     "일부만 성공하는 상태가 없습니다."),
        ("멱등성", "확정 조회가 실패해도 이체는 정확히 한 번만 집행됩니다."),
        ("위임 한도", "승인한 금액을 넘는 인출은 체인이 거부합니다. 백지수표가 아닙니다."),
        ("세션 격리", "브라우저마다 봇과 지갑이 분리됩니다. 남의 봇은 조회조차 되지 않습니다."),
    ]
    y = cm(5.4)
    for name, body in rows:
        box(s, M, y, cm(6.0), cm(1.5), fill=TINT)
        text(s, M + cm(0.6), y + cm(0.3), cm(5.2), cm(1.0), name,
             size=17, bold=True, color=BLUE)
        text(s, M + cm(6.8), y + cm(0.3), W - M - (M + cm(6.8)), cm(1.1),
             body, size=15, color=INK)
        y += cm(1.8)
    footer(s, 11)

    # ── 12. 마무리 ─────────────────────────────────────────────
    s = blank(prs)
    bg = box(s, 0, 0, W, H, fill=RGBColor(0x0B, 0x1B, 0x34), radius=False)
    bg.line.fill.background()
    text(s, M, cm(5.6), W - 2 * M, cm(2.0),
         "에이전트에게 필요한 것은 더 좋은 모델이 아니라",
         size=32, color=RGBColor(0x8F, 0xA6, 0xC4))
    text(s, M, cm(8.0), W - 2 * M, cm(2.6),
         "스스로 지불하고, 증명하고, 멈출 수 있는 능력입니다.",
         size=44, bold=True, color=WHITE)
    line = box(s, M, cm(11.6), cm(4.0), Emu(50000),
               fill=RGBColor(0x6E, 0xB4, 0xFF), radius=False)
    line.line.fill.background()
    text(s, M, cm(12.6), W - 2 * M, cm(3.0),
         "시연     googlesolana2026-ce-deploy.onrender.com\n"
         "코드     github.com/seungminkun00-hue/google-solana2026",
         size=20, color=RGBColor(0xB6, 0xC6, 0xDC))

    prs.save(OUT)
    print(f"[OK] {OUT}")
    print(f"     슬라이드 {len(prs.slides.__iter__.__self__._sldIdLst)}장")


if __name__ == "__main__":
    build()
